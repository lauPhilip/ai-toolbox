import streamlit as st
import weaviate
from weaviate.classes.init import Auth
import weaviate.classes as wvc
from weaviate.classes.query import Filter
from pypdf import PdfReader
from pptx import Presentation
import uuid
import time

# --- AUTH & ROLE GUARD ---
if st.session_state.get("authentication_status") is not True:
    st.switch_page("main.py")
    st.stop()

if str(st.session_state.get("role")).lower() != "teacher":
    st.error("Access Denied: Teacher role required.")
    st.stop()

# --- WEAVIATE CORE ---
wcd_url = st.secrets["WEAVIATE_URL"]
wcd_api_key = st.secrets["WEAVIATE_API_KEY"]

@st.cache_resource
def get_weaviate_client():
    return weaviate.connect_to_weaviate_cloud(
        cluster_url=wcd_url,
        auth_credentials=Auth.api_key(wcd_api_key),
    )

client = get_weaviate_client()
collection = client.collections.get("CourseBotMemory")
course_registry = client.collections.get("CourseRegistry")

# --- INITIALIZE RESET KEYS ---
if "deploy_key" not in st.session_state:
    st.session_state.deploy_key = 0

def reset_deploy_form():
    st.session_state.deploy_key += 1
    st.rerun()

# --- UTILITY: FILE EXTRACTION ---
def extract_text(file):
    try:
        if file.name.endswith(".pdf"):
            reader = PdfReader(file)
            return " ".join([page.extract_text() for page in reader.pages])
        elif file.name.endswith(".pptx"):
            prs = Presentation(file)
            return " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except Exception as e:
        st.error(f"Failed to read {file.name}: {e}")
    return ""

st.title(f"👨‍🏫 {st.session_state['name']}'s course-bot Dashboard")
tab_manage, tab_upload = st.tabs(["📚 Managed CourseBots", "🤖 Deploy New Bot"])

# --- TAB 1: DEPLOY NEW BOT ---
with tab_upload:
    st.subheader("Deploy a Specialized CourseBot")
    
    user_email = st.session_state.get("email")
    registry_resp = course_registry.query.fetch_objects(
        filters=Filter.by_property("responsible_email").equal(user_email) | 
                Filter.by_property("teacher_emails").contains_any([user_email]),
        return_properties=["course_name", "course_id"],
        limit=100
    )
    
    course_options = {f"{obj.properties['course_name']} ({obj.properties['course_id']})": obj.properties['course_id'] 
                      for obj in registry_resp.objects}

    if not course_options:
        st.warning("⚠️ No courses found in the Registrar.")
    else:
        # Form-style container with unique keys derived from deploy_key for easy resetting
        with st.container(border=True):
            sel_course_label = st.selectbox("Assign to Course", options=list(course_options.keys()), key=f"course_sel_{st.session_state.deploy_key}")
            sel_course_id = course_options[sel_course_label]
            
            b_name = st.text_input("CourseBot Name", placeholder="e.g., Exam Prep", key=f"bot_name_{st.session_state.deploy_key}")
            p_level = st.selectbox("Level", ["Bachelor", "Master"], key=f"lvl_{st.session_state.deploy_key}")

            st.write("---")
            sys_p = st.text_area("System Prompt", value="You are a professional academic assistant...", height=200, key=f"prompt_{st.session_state.deploy_key}")
            temp = st.slider("Temperature (Creativity)", 0.0, 1.0, 0.2, 0.1, key=f"temp_{st.session_state.deploy_key}")

            up_files = st.file_uploader("Upload Knowledge Base", accept_multiple_files=True, key=f"uploader_{st.session_state.deploy_key}")
            
            if st.button("🚀 Deploy CourseBot", type="primary", use_container_width=True):
                if b_name and up_files:
                    with st.spinner(f"Deploying {b_name}..."):
                        for file in up_files:
                            text = extract_text(file)
                            chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
                            for chunk in chunks:
                                collection.data.insert(properties={
                                    "doc_title": file.name, "chunk": chunk, "bot_name": b_name,
                                    "course_id": sel_course_id, "course_name": sel_course_label,
                                    "course_administrator": user_email, "program": p_level,
                                    "system_prompt": sys_p, "temperature": float(temp)
                                })
                    st.success(f"'{b_name}' deployed!")
                    time.sleep(1)
                    reset_deploy_form()
                else:
                    st.warning("Bot Name and Files are required.")

# --- TAB 2: MANAGE & EDIT ---
with tab_manage:
    st.subheader("Manage Active CourseBots")
    
    user_email = st.session_state.get("email")
    results = collection.query.fetch_objects(
        filters=Filter.by_property("course_administrator").equal(user_email),
        return_properties=["bot_name", "course_name"], limit=1000 
    )

    organized = {}
    for obj in results.objects:
        c_name = obj.properties['course_name']
        b_name = obj.properties.get('bot_name') or "Standard Bot"
        if c_name not in organized: organized[c_name] = []
        if b_name not in organized[c_name]: organized[c_name].append(b_name)
    
    if organized:
        sel_c = st.selectbox("Filter by Course", options=list(organized.keys()), key="m_c_sel")
        sel_b = st.selectbox("Select CourseBot", options=organized[sel_c], key="m_b_sel")
        
        st.write("---")
        st.markdown(f"### ⚙️ Configuration: **{sel_b}**")
        
        # 1. SETTINGS MANAGEMENT
        config = collection.query.fetch_objects(
            filters=Filter.by_property("course_name").equal(sel_c) & Filter.by_property("bot_name").equal(sel_b),
            return_properties=["system_prompt", "temperature"], limit=1
        )
        
        if config.objects:
            obj = config.objects[0]
            edit_p = st.text_area("System Prompt", value=obj.properties.get("system_prompt"), height=200, key="m_p_area")
            edit_t = st.slider("Temperature", 0.0, 1.0, float(obj.properties.get("temperature", 0.2)), 0.1, key="m_t_slider")
            
            if st.button("💾 Save Changes", type="primary", key="m_save_btn"):
                targets = collection.query.fetch_objects(
                    filters=Filter.by_property("course_name").equal(sel_c) & Filter.by_property("bot_name").equal(sel_b),
                    return_properties=[], limit=10000
                )
                for target in targets.objects:
                    collection.data.update(uuid=target.uuid, properties={"system_prompt": edit_p, "temperature": edit_t})
                st.success("Tuning complete.")
                st.rerun()

        # 2. FILE MANAGEMENT (Edit/Delete existing files)
        st.write("---")
        st.markdown(f"### 📄 Knowledge Assets for **{sel_b}**")
        
        file_results = collection.query.fetch_objects(
            filters=Filter.by_property("course_name").equal(sel_c) & Filter.by_property("bot_name").equal(sel_b),
            return_properties=["doc_title"], limit=1000
        )
        unique_files = sorted(list(set([o.properties['doc_title'] for o in file_results.objects])))

        if not unique_files:
            st.info("No files currently assigned to this bot.")
        else:
            for f_name in unique_files:
                col_txt, col_del = st.columns([0.85, 0.15])
                col_txt.markdown(f"📄 `{f_name}`")
                if col_del.button("🗑️", key=f"del_file_{f_name}_{sel_b}"):
                    collection.data.delete_many(
                        where=Filter.by_property("doc_title").equal(f_name) & 
                              Filter.by_property("bot_name").equal(sel_b) &
                              Filter.by_property("course_name").equal(sel_c)
                    )
                    st.toast(f"Purged {f_name}")
                    time.sleep(0.5)
                    st.rerun()

        # 3. ADD NEW CONTENT
        st.markdown("#### ➕ Add Knowledge")
        new_up = st.file_uploader("Upload additional files", accept_multiple_files=True, key="m_file_add")
        if st.button("🚀 Process & Add", key="m_add_btn"):
            if new_up:
                with st.spinner("Integrating..."):
                    for file in new_up:
                        text = extract_text(file)
                        chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
                        # Using existing edit_p and edit_t to stay consistent
                        for chunk in chunks:
                            collection.data.insert(properties={
                                "doc_title": file.name, "chunk": chunk, "bot_name": sel_b,
                                "course_name": sel_c, "course_administrator": user_email,
                                "system_prompt": edit_p, "temperature": edit_t
                            })
                st.rerun()

        # 4. DANGER ZONE
        with st.expander("⚠️ Danger Zone"):
            if st.button(f"🔥 Wipe '{sel_b}' Permanently", type="primary", use_container_width=True):
                collection.data.delete_many(
                    where=Filter.by_property("course_name").equal(sel_c) & Filter.by_property("bot_name").equal(sel_b)
                )
                st.rerun()
    else:
        st.info("No active CourseBots found.")